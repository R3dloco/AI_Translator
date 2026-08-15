import os
import io
import email
from typing import Optional, Dict, List
import pandas as pd
import pypdf
import docx
import pptx
from bs4 import BeautifulSoup
import ebooklib
from ebooklib import epub
import extract_msg

class DocumentExtractor:
    """Extracts plain text from various document formats."""
    
    SUPPORTED_EXTENSIONS = {
        # Text & Markdown
        ".txt", ".md", ".rtf", ".html", ".htm", ".json", ".yaml", ".yml", ".xml", ".log",
        ".rst", ".tex", ".ini", ".conf", ".cfg", ".env", ".properties", ".toml", ".sql",
        ".srt", ".vtt",
        # Source Code
        ".py", ".js", ".ts", ".css", ".cpp", ".c", ".h", ".hpp", ".hxx", ".java",
        ".cs", ".go", ".php", ".rb", ".swift", ".kt", ".scala", ".r", ".sh", ".bash",
        ".jsx", ".tsx", ".vue",
        # PDF
        ".pdf",
        # Word
        ".docx",
        # PowerPoint
        ".pptx",
        # Excel & CSV
        ".xlsx", ".xls", ".csv",
        # E-books
        ".epub",
        # Emails
        ".eml", ".msg"
    }

    @classmethod
    def is_supported(cls, filename: str) -> bool:
        ext = os.path.splitext(filename)[1].lower()
        return ext in cls.SUPPORTED_EXTENSIONS

    @classmethod
    def get_supported_extensions(cls) -> List[str]:
        return sorted(list(cls.SUPPORTED_EXTENSIONS))

    def extract_text(self, file_path_or_bytes, filename: str) -> str:
        """
        Extract text from file_path or file bytes based on file extension.
        `file_path_or_bytes` can be a str path or bytes/BytesIO stream.
        """
        ext = os.path.splitext(filename)[1].lower()
        
        if isinstance(file_path_or_bytes, (str, bytes)):
            if isinstance(file_path_or_bytes, str) and os.path.isfile(file_path_or_bytes):
                with open(file_path_or_bytes, "rb") as f:
                    content_bytes = f.read()
            elif isinstance(file_path_or_bytes, bytes):
                content_bytes = file_path_or_bytes
            else:
                raise FileNotFoundError(f"File not found: {file_path_or_bytes}")
        elif hasattr(file_path_or_bytes, "read"):
            content_bytes = file_path_or_bytes.read()
        else:
            raise ValueError("Unsupported input type for document extraction.")

        if ext == ".pdf":
            return self._extract_pdf(content_bytes)
        elif ext == ".docx":
            return self._extract_docx(content_bytes)
        elif ext == ".pptx":
            return self._extract_pptx(content_bytes)
        elif ext in (".xlsx", ".xls"):
            return self._extract_excel(content_bytes)
        elif ext == ".csv":
            return self._extract_csv(content_bytes)
        elif ext == ".epub":
            return self._extract_epub(content_bytes)
        elif ext == ".eml":
            return self._extract_eml(content_bytes)
        elif ext == ".msg":
            return self._extract_msg(content_bytes)
        elif ext in (".html", ".htm"):
            return self._extract_html(content_bytes)
        else:
            # Fallback for plain text, markdown, json, code, etc.
            return self._extract_plain_text(content_bytes)

    def _extract_pdf(self, content_bytes: bytes) -> str:
        pdf_file = io.BytesIO(content_bytes)
        reader = pypdf.PdfReader(pdf_file)
        text_parts = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            # If native text is missing or extremely short, attempt OCR on scanned page
            if len(page_text.strip()) < 10:
                try:
                    import pdf2image
                    import pytesseract
                    images = pdf2image.convert_from_bytes(content_bytes, first_page=i+1, last_page=i+1)
                    if images:
                        project_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
                        local_tessdata = os.path.join(project_dir, "tessdata")
                        config = f'--tessdata-dir "{local_tessdata}"' if os.path.isfile(os.path.join(local_tessdata, "eng.traineddata")) else ""
                        ocr_text = pytesseract.image_to_string(images[0], config=config).strip() if config else pytesseract.image_to_string(images[0]).strip()
                        if ocr_text:
                            page_text = f"[OCR Extracted]\n{ocr_text}"
                except Exception:
                    pass

            if page_text.strip():
                text_parts.append(f"--- Page {i + 1} ---\n{page_text.strip()}")
        return "\n\n".join(text_parts)

    def _extract_docx(self, content_bytes: bytes) -> str:
        docx_file = io.BytesIO(content_bytes)
        doc = docx.Document(docx_file)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        return "\n".join(paragraphs)

    def _extract_pptx(self, content_bytes: bytes) -> str:
        pptx_file = io.BytesIO(content_bytes)
        prs = pptx.Presentation(pptx_file)
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_parts = [f"--- Slide {i + 1} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if len(slide_parts) > 1:
                slides_text.append("\n".join(slide_parts))
        return "\n\n".join(slides_text)

    def _extract_excel(self, content_bytes: bytes) -> str:
        excel_file = io.BytesIO(content_bytes)
        excel = pd.read_excel(excel_file, sheet_name=None)
        sheets_text = []
        for sheet_name, df in excel.items():
            sheets_text.append(f"=== Sheet: {sheet_name} ===")
            sheets_text.append(df.to_csv(index=False))
        return "\n\n".join(sheets_text)

    def _extract_csv(self, content_bytes: bytes) -> str:
        text = self._extract_plain_text(content_bytes)
        return text

    def _extract_epub(self, content_bytes: bytes) -> str:
        epub_file = io.BytesIO(content_bytes)
        book = epub.read_epub(epub_file)
        chapters = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), 'html.parser')
            text = soup.get_text(separator='\n').strip()
            if text:
                chapters.append(text)
        return "\n\n".join(chapters)

    def _extract_eml(self, content_bytes: bytes) -> str:
        msg = email.message_from_bytes(content_bytes)
        parts = []
        subject = msg.get('subject', '')
        sender = msg.get('from', '')
        date = msg.get('date', '')
        parts.append(f"Subject: {subject}\nFrom: {sender}\nDate: {date}\n")
        
        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                if content_type == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        parts.append(payload.decode('utf-8', errors='replace'))
                elif content_type == "text/html":
                    payload = part.get_payload(decode=True)
                    if payload:
                        soup = BeautifulSoup(payload.decode('utf-8', errors='replace'), 'html.parser')
                        parts.append(soup.get_text(separator='\n'))
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                parts.append(payload.decode('utf-8', errors='replace'))
        return "\n\n".join(parts)

    def _extract_msg(self, content_bytes: bytes) -> str:
        msg = extract_msg.Message(io.BytesIO(content_bytes))
        body = msg.body or ""
        header = f"Subject: {msg.subject}\nFrom: {msg.sender}\nDate: {msg.date}\n"
        return f"{header}\n{body}"

    def _extract_html(self, content_bytes: bytes) -> str:
        html_str = self._extract_plain_text(content_bytes)
        soup = BeautifulSoup(html_str, 'html.parser')
        return soup.get_text(separator='\n')

    def _extract_plain_text(self, content_bytes: bytes) -> str:
        for encoding in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']:
            try:
                return content_bytes.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content_bytes.decode('utf-8', errors='replace')
